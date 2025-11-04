from flask import Flask, render_template, request, send_from_directory, jsonify
import os
from werkzeug.utils import secure_filename
from docx import Document
import fitz  # for PDF
from pptx import Presentation
from PIL import Image
import qrcode

app = Flask(__name__)
UPLOAD_FOLDER = 'uploads'
CONVERTED_FOLDER = 'converted'
COMBINED_FOLDER = 'combined'

for folder in [UPLOAD_FOLDER, CONVERTED_FOLDER, COMBINED_FOLDER]:
    os.makedirs(folder, exist_ok=True)

@app.route('/')
def home():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    files = request.files.getlist('files')
    converted_files = []
    for file in files:
        filename = secure_filename(file.filename)
        base, ext = os.path.splitext(filename)
        filepath = os.path.join(UPLOAD_FOLDER, filename)
        file.save(filepath)
        output_path = os.path.join(CONVERTED_FOLDER, base + '.docx')

        doc = Document()
        if ext.lower() == '.pdf':
            pdf = fitz.open(filepath)
            for page in pdf:
                text = page.get_text()
                doc.add_paragraph(text)
            pdf.close()
        elif ext.lower() in ['.jpg', '.jpeg', '.png']:
            doc.add_paragraph(f"Image inserted: {filename}")
        elif ext.lower() == '.pptx':
            ppt = Presentation(filepath)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        doc.add_paragraph(shape.text)
        else:
            continue
        doc.save(output_path)
        converted_files.append(os.path.basename(output_path))
    return jsonify({'converted': converted_files})

@app.route('/converted')
def list_converted():
    files = os.listdir(CONVERTED_FOLDER)
    return jsonify(files)

@app.route('/download/<folder>/<filename>')
def download_file(folder, filename):
    return send_from_directory(folder, filename, as_attachment=True)

@app.route('/combine', methods=['POST'])
def combine_files():
    files = request.json.get('files', [])
    combined_doc = Document()
    output_path = os.path.join(COMBINED_FOLDER, 'Combined_Wordify.docx')
    for fname in files:
        path = os.path.join(CONVERTED_FOLDER, fname)
        if os.path.exists(path):
            subdoc = Document(path)
            for element in subdoc.element.body:
                combined_doc.element.body.append(element)
    combined_doc.save(output_path)
    return jsonify({'combined': 'Combined_Wordify.docx'})

@app.route('/generate_qr', methods=['POST'])
def generate_qr():
    text = request.json.get('text', '')
    if not text:
        return jsonify({'error': 'No text provided'})
    img = qrcode.make(text)
    qr_path = os.path.join('static', 'qr.png')
    os.makedirs('static', exist_ok=True)
    img.save(qr_path)
    return jsonify({'qr_path': qr_path})

if __name__ == '__main__':
    app.run(debug=True)
